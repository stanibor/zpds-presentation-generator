from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1 - Title Slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)

# Set a custom background image for Slide 1
background_image = "path/to/background_image.jpg"
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background color
slide.background.fill.background()

# Title
title = slide.shapes.title
title.text = "Academic Presentation"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text_frame.margin_left = Inches(1)
title.text_frame.margin_right = Inches(1)
title.text_frame.margin_bottom = Inches(1)

# Subtitle
subtitle = slide.placeholders[1]
subtitle.text = "Your Name"
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
subtitle.text_frame.margin_left = Inches(1)
subtitle.text_frame.margin_right = Inches(1)
subtitle.text_frame.margin_bottom = Inches(0.5)
subtitle.text_frame.margin_top = Inches(0)

# Slide 2 - Content Slide
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)

# Set a solid background color for Slide 2
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(41, 128, 185)  # Blue background color

# Slide Title
title = slide.shapes.title
title.text = "Introduction"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text_frame.margin_left = Inches(1)
title.text_frame.margin_right = Inches(1)
title.text_frame.margin_bottom = Inches(0.5)
title.text_frame.margin_top = Inches(0.5)
title.text_frame.word_wrap = True

# Slide Content
content = slide.placeholders[1]
content.text = "Lorem ipsum dolor sit amet, consectetur adipiscing elit."
content.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
content.text_frame.margin_left = Inches(1)
content.text_frame.margin_right = Inches(1)
content.text_frame.margin_bottom = Inches(1)
content.text_frame.word_wrap = True

# Add a sleek design shape to Slide 2
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(8), Inches(0.5))

# Customize the shape properties
shape = slide.shapes[-1]
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill color

outline = shape.line
outline.width = Pt(2)
outline.color.rgb = RGBColor(41, 128, 185)  # Blue outline color

# Save the PowerPoint presentation
presentation.save("sophisticated_presentation.pptx")