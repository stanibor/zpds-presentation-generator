from pathlib import Path

import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from audio_functions import get_pretrained_tts_models, synthesize_presentation_notes, extract_presentation_notes, \
    save_waveforms, annotate_presentation_with_spoken_notes
from defined_functions import get_presenation_data


def main():
    st.markdown("<h1 style='text-align: center; color: White;background-color:#e84343'>Presentation Creator</h1>", unsafe_allow_html=True)
    st.markdown("<h3 style='text-align: center; color: Black;'>Drop in The required Inputs and we will do  the rest.</h3>", unsafe_allow_html=True)
    st.sidebar.header("What is this Project about?")
    st.sidebar.text("It a Web app that would help the user create presentations and present them.")
    st.sidebar.header("What tools where used to make this?")
    st.sidebar.text("Chat GPT3 and more.")

    n_of_slides = st.number_input("Input number of core slides", 1, 20)
    words = st.text_input("Input context words, semicolon separated")
    list_of_words = words.split(";")

    if 'dwn_avail' not in st.session_state:
        st.session_state.dwn_avail = 0
    if 'dwn_tts_avail' not in st.session_state:
        st.session_state.dwn_tts_avail = 0


    if st.button("Generate slides"):
        with st.spinner('Wait for it...'):
            st.session_state.dwn_avail = 1
            data = get_presenation_data(list_of_words, n_of_slides)
            with open('slide_notes.txt', 'w', encoding="utf8") as f:
                f.write(data['introduction_speech'])
                for slide_data in data['slides']:
                    f.write(slide_data['speech'])
                    f.write('\n')
                f.write(data['summary'])

            pres = Presentation()

            #Slide Intro

            slide = pres.slides.add_slide(pres.slide_layouts[2])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(173, 216, 230)
            title_shape = slide.shapes.title
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(96, 147, 172)
            title_shape.text = data['title']
            left = top = Inches(2)
            width = height = Inches(3)

            notes_slide = slide.notes_slide  # Get slide notes object
            notes_slide.notes_text_frame.text = data["introduction_speech"]

            #Slides

            for slide_data in data['slides']:
                slide_layout = pres.slide_layouts[8]  # Use layout with title and content
                slide = pres.slides.add_slide(slide_layout)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(173, 216, 230)
                # Set the slide title
                title = slide.shapes.title
                title_shape.fill.solid()
                title_shape.fill.fore_color.rgb = RGBColor(96, 147, 172)
                title.text = slide_data['title']
                # Add a text field
                content = slide.placeholders[2]  # Use the second placeholder for content
                content.text = slide_data['text']
                # Add a picture
                if len(slide_data['img'][1]):
                    picture_path = slide_data['img'][1]  # Replace with the actual path to your picture
                    img_placeholder = slide.placeholders[1]
                    img_placeholder.insert_picture(picture_path)

                # Add speech into slide notes
                notes_slide = slide.notes_slide  # Get slide notes object
                notes_slide.notes_text_frame.text = slide_data['speech']  # Get speech text into notes

            slide = pres.slides.add_slide(pres.slide_layouts[2])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(173, 216, 230)
            title_shape = slide.shapes.title
            title_shape.text = 'Thank you'
            left = top = Inches(2)
            width = height = Inches(3)

            notes_slide = slide.notes_slide  # Get slide notes object
            notes_slide.notes_text_frame.text = data['summary']  # Get speech text into notes

            pres.save('presentation.pptx')
        st.success('Done!')
    
    if st.session_state.dwn_avail == 1:
        with open('presentation.pptx', "rb") as file:  # , encoding="ISO-8859-1"
            btn = st.download_button(
                    label="Download presentation .pptx",
                    data=file,
                    file_name='presentation.pptx',
                    disabled=(st.session_state.dwn_avail == 0)
                )
            # presentation = file

    presentation_dl = st.file_uploader("Upload presentation .pptx", type=["pptx"])
    presentation = Presentation("presentation.pptx") if Path('presentation.pptx').exists() else None

    if presentation_dl is not None:
        presentation = Presentation(presentation_dl)

    # Not sure how it works maybe it has to be cached once or something
    tacotron_model, hifi_gan = get_pretrained_tts_models()

    if st.button("Generate a spoken presentation"):
        if presentation is None:
            st.write("Files are not provided!")
        else:
            progress_text = "Generating speech in progress. Please wait."
            my_bar = st.progress(0, text=progress_text)
            spoken_notes = []
            presentation_len = len(presentation.slides)
            slide_notes = extract_presentation_notes(presentation)
            for slide_wav in synthesize_presentation_notes(slide_notes, tacotron_model, hifi_gan):
                i = len(spoken_notes) + 1
                progress_text = f"Generating speech in progress. Please wait. (Slide {i}/{presentation_len})"
                spoken_notes.append(slide_wav)
                my_bar.progress(i / presentation_len, text=progress_text)
            progress_text = f"Generating speech completed."
            my_bar.progress(1.0, text=progress_text)
            with st.spinner('Adding audio tracks...'):
                # proper audio tracks can be added to each slide of the presentation
                audio_tracks = save_waveforms(spoken_notes)
                annotate_presentation_with_spoken_notes(presentation, audio_tracks)

                presentation.save("spoken_presentation.pptx")

                st.session_state.dwn_tts_avail = 1

    if st.session_state.dwn_tts_avail == 1:
        with open("spoken_presentation.pptx", "rb") as file:
            st.download_button("Download spoken presentation .pptx",
                               data=file,
                               file_name='presentation.pptx',
                               disabled=(st.session_state.dwn_tts_avail == 0))


if __name__ == '__main__':
    main()
