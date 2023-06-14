from pptx import Presentation
import numpy as np
import torchaudio
import torch
from pathlib import Path
from os import PathLike
from speechbrain.pretrained import Tacotron2, HIFIGAN

from typing import List

DEFAULT_SAMPLING_RATE = 22050

def get_pretrained_tacotron(tmpdir: Path = Path("")):
    """
    Get pretrained tacotron model from SpeechBrain library and save it to tmpdir
    :param tmpdir: directory to save the model
    :return: tacotron model
    """
    # prepare temporary directories for model weights and metadata
    tmpdir = Path(tmpdir)
    tmpdir_vocoder = tmpdir / "vocoder"
    tacotron_model = Tacotron2.from_hparams(source="speechbrain/tts-tacotron2-ljspeech", savedir=tmpdir_vocoder)
    return tacotron_model


def get_pretrained_hifi_gan(tmpdir: Path = Path("")):
    """
    Get pretrained hifi-gan model from SpeechBrain library and save it to tmpdir
    :param tmpdir: directory to save the model
    :return: hifi-gan model
    """
    # prepare temporary directories for model weights and metadata
    tmpdir = Path(tmpdir)
    tmpdir_tts = tmpdir / "tts"
    hifi_gan = HIFIGAN.from_hparams(source="speechbrain/tts-hifigan-ljspeech", savedir=tmpdir_tts)
    return hifi_gan


def get_pretrained_tts_models(tmpdir: Path = Path("")):
    """
    Get pretrained tacotron and hifi-gan models from SpeechBrain library and save it to tmpdir
    :param tmpdir: directory to save the models
    :return: tuple of tacotron and hifi-gan models
    """
    tacotron_model = get_pretrained_tacotron(tmpdir)
    hifi_gan = get_pretrained_hifi_gan(tmpdir)
    return tacotron_model, hifi_gan


def extract_presentation_notes(presentation: Presentation) -> List[str]:
    """
    Extract notes from presentation slides
    :param presentation: presentation object
    :return: list of notes strings for each slide
    """
    presentation_notes = [slide.notes_slide.notes_text_frame.text for slide in presentation.slides]
    return presentation_notes


def synthesize_slide_notes(slide_notes: List[str], tacotron_model, hifi_gan) -> List[torch.Tensor]:
    """
    Synthesize speech for a list of notes. Each note is a sentence.
    :param slide_notes: list of notes
    :param tacotron_model: model to synthesize speech
    :param hifi_gan: model to convert mel spectrogram to waveform
    :return: list of waveform tensors for each note
    """
    wavs = []
    for note in slide_notes:
        mel_outputs, mel_lengths, alignments = tacotron_model.encode_batch([note])
        waveform = hifi_gan.decode_batch(mel_outputs)[0]
        wavs.append(waveform)
    return wavs


def consolidate_speech(wavs: List[torch.Tensor],
                      pause_time=0.5,
                      sampling_rate = DEFAULT_SAMPLING_RATE):
    """
    Concatenate speech waveforms with pauses between them. Pauses are added to make the speech more natural.
    :param wavs: list of speech waveform tensors to concatenate
    :param pause_time: duration of pauses in seconds
    :param sampling_rate: sampling rate of the waveforms
    :return: concatenated waveform tensor
    """
    pause = wavs[0].new_zeros(1, int(pause_time * sampling_rate))
    wavs_pauses = [item for wav in wavs for item in [wav, pause]]
    big_wav = torch.cat(wavs_pauses, dim=-1)
    return big_wav


def synthesize_presentation_notes(presentation_notes: List[str], tacotron_model, hifi_gan) -> List[torch.Tensor]:
    """
    Synthesize speech for a list of notes. Each note is a full speech for a slide.

    :param presentation_notes: list of notes
    :param tacotron_model: model to synthesize speech
    :param hifi_gan: model to convert mel spectrogram to waveform
    :return: list of waveform tensors for each slide
    """
    presentation_notes = [t.replace('"', "") for t in presentation_notes]
    # wavs = []
    for slide_notes in presentation_notes:
        slide_sentences = slide_notes.split(".")
        slide_sentences = [s for s in slide_sentences if len(s) > 0]
        slide_wavs = synthesize_slide_notes(slide_sentences, tacotron_model, hifi_gan)
        slide_wav = consolidate_speech(slide_wavs)
        # wavs.append(slide_wav)
        yield slide_wav
    # return wavs


def save_full_speech(wavs: List[torch.Tensor],
                     output_path: Path = Path("full_speech.wav"),
                     pause_time=0.5,
                     sampling_rate = DEFAULT_SAMPLING_RATE):
    """
    Concatenate speech waveforms with pauses between them and save the result to a file.
    Pauses are added to make the speech more natural.
    :param wavs: list of speech waveform tensors to concatenate
    :param output_path: path to save the result
    :param pause_time: duration of pauses in seconds
    :param sampling_rate: sampling rate of the waveforms
    """
    big_wav = consolidate_speech(wavs, pause_time, sampling_rate)
    torchaudio.save(output_path, big_wav, sampling_rate)


def save_waveforms(wavs: List[torch.Tensor], output_dir=Path(""), sampling_rate=DEFAULT_SAMPLING_RATE):
    """
    Save a list of waveforms to a directory. Each waveform is saved to a separate file.
    The file names are S0.wav, S1.wav, etc.
    :param wavs: list of waveforms
    :param output_dir: directory to save the waveforms
    :param sampling_rate: sampling rate of the waveforms
    :return: list of paths to the saved waveforms
    """
    output_dir = Path(output_dir)
    audio_tracks = []
    for i, wav in enumerate(wavs):
        audio_track = output_dir / f"S{i}.wav"
        torchaudio.save(str(audio_track), wav, sampling_rate)
        audio_tracks.append(str(audio_track))

    return audio_tracks


from lxml import etree
def annotate_presentation_with_spoken_notes(presentation: Presentation, audio_tracks: List[PathLike]) -> Presentation:
    """
    Add audio tracks to presentation slides. The audio tracks are added as movies to the slides.
    The audio tracks are played automatically when the slide is shown.
    :param presentation: presentation object
    :param audio_tracks: list of paths to audio tracks
    :return: presentation object with audio tracks added
    """
    assert len(audio_tracks) <= len(presentation.slides)
    for i, audio_track in enumerate(audio_tracks):
        slide = presentation.slides[i]
        movie = slide.shapes.add_movie(audio_track, 0, 0, 100, 100)
        tree = movie._element.getparent().getparent().getnext().getnext()
        timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
        timing.set('delay', '0')

    return presentation

